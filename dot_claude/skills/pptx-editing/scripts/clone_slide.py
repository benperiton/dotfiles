#!/usr/bin/env python
"""Duplicate a slide faithfully -- the operation python-pptx has no API for.

Why this exists: matching a deck's look means cloning a real slide, not building
from a (often near-empty) layout. python-pptx can't duplicate a slide, so people
hand-edit the zip and corrupt it. This does the copy at the XML level but lets
python-pptx own the dangerous plumbing ([Content_Types], .rels, sldIdLst), so the
output always opens.

What it does:
  1. Adds a new slide using the SOURCE slide's own layout.
  2. Strips the placeholders add_slide() inherited from that layout.
  3. Deep-copies every shape element from the source slide's spTree.
  4. Re-creates the relationships those shapes reference (images, media,
     hyperlinks, OLE) on the new slide part and rewrites the r:id/r:embed/r:link
     attributes to the new ids. Image parts are shared, not duplicated.
  5. Optionally moves the new slide after a given 1-based position.

It does NOT edit text -- clone first, then edit the copy with python-pptx and
verify with render.sh + validate.py.

Usage:
  clone_slide.py DECK.pptx --source N [--after M] [--out OUT.pptx] [--times K]

  --source N   1-based slide number to clone (required)
  --after  M   place the clone after slide M (1-based). Omit = append at end.
  --times  K   clone K copies (default 1); with --after they land in order.
  --out        write here instead of editing DECK.pptx in place.

Run with the dedicated venv:
  ~/.venvs/pptx/bin/python clone_slide.py deck.pptx --source 7 --after 7
"""
import argparse
import copy

from pptx import Presentation
from pptx.oxml.ns import qn

RELS_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_KEEP = {qn("p:nvGrpSpPr"), qn("p:grpSpPr")}


def _rel_attrs(elem):
    """Yield (element, attr_key) for every relationship-namespaced attribute."""
    for node in elem.iter():
        for key in node.attrib:
            if key.startswith("{" + RELS_NS + "}"):
                yield node, key


def clone_slide(prs, source_index):
    """Clone the 1-based source_index slide; return the new Slide (appended last)."""
    source = prs.slides[source_index - 1]
    dest = prs.slides.add_slide(source.slide_layout)

    dest_tree = dest.shapes._spTree
    for child in list(dest_tree):
        if child.tag not in _KEEP:
            dest_tree.remove(child)

    src_tree = source.shapes._spTree
    for child in list(src_tree):
        if child.tag not in _KEEP:
            dest_tree.append(copy.deepcopy(child))

    # Re-create referenced relationships on the new part and remap the ids.
    id_map = {}
    for node, key in _rel_attrs(dest_tree):
        old_rid = node.get(key)
        if old_rid in id_map:
            node.set(key, id_map[old_rid])
            continue
        rel = source.part.rels[old_rid]
        if rel.is_external:
            new_rid = dest.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_rid = dest.part.relate_to(rel.target_part, rel.reltype)
        id_map[old_rid] = new_rid
        node.set(key, new_rid)

    return dest


def move_after(prs, from_index_1based, after_1based):
    """Move the slide currently at from_index (1-based) to follow after_1based."""
    sld_id_lst = prs.slides._sldIdLst
    nodes = list(sld_id_lst)
    moving = nodes[from_index_1based - 1]
    sld_id_lst.remove(moving)
    sld_id_lst.insert(after_1based, moving)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("deck")
    ap.add_argument("--source", type=int, required=True)
    ap.add_argument("--after", type=int, default=None)
    ap.add_argument("--times", type=int, default=1)
    ap.add_argument("--out", default=None)
    args = ap.parse_args()

    prs = Presentation(args.deck)
    n = len(prs.slides._sldIdLst)
    if not 1 <= args.source <= n:
        ap.error(f"--source {args.source} out of range (deck has {n} slides)")

    insert_at = args.after  # advances so multiple clones stay in order
    for _ in range(args.times):
        clone_slide(prs, args.source)
        if insert_at is not None:
            move_after(prs, len(prs.slides._sldIdLst), insert_at)
            insert_at += 1

    out = args.out or args.deck
    prs.save(out)
    print(f"cloned slide {args.source} x{args.times} -> {out} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
