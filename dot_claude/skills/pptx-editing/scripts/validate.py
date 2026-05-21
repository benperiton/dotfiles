#!/usr/bin/env python
"""Integrity check for a .pptx after editing.

Catches the silent-corruption class of bug: a deck that python-pptx wrote but
that PowerPoint would refuse to open (or "repair", dropping content). Run this
after every edit, and compare slide count against what you expected.

Checks:
  1. The file is a valid zip and every entry passes CRC (zipfile.testzip).
  2. [Content_Types].xml is present.
  3. python-pptx can reopen it (parses presentation.xml + every slide part).
  4. Every relationship referenced by each slide actually resolves to a part.

Exit code is 0 only if all checks pass. Prints the slide count and per-slide
shape count so you can eyeball that nothing was lost.

Usage:
  validate.py DECK.pptx [--expect-slides N]

Run with the dedicated venv:
  ~/.venvs/pptx/bin/python validate.py deck.pptx --expect-slides 21
"""
import argparse
import sys
import zipfile

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("deck")
    ap.add_argument("--expect-slides", type=int, default=None)
    args = ap.parse_args()

    problems = []

    # 1 + 2: zip integrity and required part.
    try:
        with zipfile.ZipFile(args.deck) as zf:
            bad = zf.testzip()
            if bad is not None:
                problems.append(f"corrupt zip entry: {bad}")
            if "[Content_Types].xml" not in zf.namelist():
                problems.append("missing [Content_Types].xml")
    except zipfile.BadZipFile as exc:
        print(f"FAIL: not a valid zip/pptx: {exc}", file=sys.stderr)
        return 1

    # 3: python-pptx can reopen and parse every slide.
    try:
        prs = Presentation(args.deck)
        slides = list(prs.slides)
    except Exception as exc:  # noqa: BLE001 - reopen failure is the signal
        print(f"FAIL: python-pptx could not reopen: {exc}", file=sys.stderr)
        return 1

    print(f"slides: {len(slides)}")
    for i, slide in enumerate(slides, 1):
        shapes = list(slide.shapes)
        print(f"  slide {i:>2}: {len(shapes)} shapes")
        # 4: every relationship target resolves.
        for rid, rel in slide.part.rels.items():
            if rel.is_external:
                continue
            try:
                _ = rel.target_part
            except Exception as exc:  # noqa: BLE001
                problems.append(f"slide {i} rel {rid} ({rel.reltype}) dangling: {exc}")
        # A slide must reference exactly one layout.
        layouts = [r for r in slide.part.rels.values() if r.reltype == RT.SLIDE_LAYOUT]
        if len(layouts) != 1:
            problems.append(f"slide {i} references {len(layouts)} layouts (expected 1)")

    if args.expect_slides is not None and len(slides) != args.expect_slides:
        problems.append(f"slide count {len(slides)} != expected {args.expect_slides}")

    if problems:
        print("\nFAIL:", file=sys.stderr)
        for p in problems:
            print(f"  - {p}", file=sys.stderr)
        return 1

    print("\nOK: deck is structurally sound")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
